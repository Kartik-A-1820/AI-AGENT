from langchain.tools import tool
from pydantic import BaseModel, Field
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches
import os
from image_generation_tools import generate_image_from_prompt

class DocxContent(BaseModel):
    text: str = Field(..., description="A paragraph of text for the document.")
    image_path: Optional[str] = Field(None, description="The local path to an image to include with the paragraph.")


AGENT_DIR = "agent_working"

class SlideData(BaseModel):
    title: str = Field(..., description="The title of the slide.")
    content: List[str] = Field(..., description="A list of bullet points for the slide's content.")
    image_path: Optional[str] = Field(
        None, description="Local path to an image for the slide."
    )
    image_prompt: Optional[str] = Field(
        None, description="Text prompt to auto-generate an image when no local path is given."
    )

@tool
def create_pptx_presentation(
    slides_data: List[SlideData],
    output_filename: str = "presentation.pptx"
) -> str:
    """
    Generates a PowerPoint presentation and inserts images per slide.
    If `image_prompt` is provided without `image_path`, it auto-generates
    the image via `generate_image_from_prompt.invoke(...)`.
    """
    # Ensure output directory exists
    os.makedirs(AGENT_DIR, exist_ok=True)
    prs = Presentation()

    for idx, slide_info in enumerate(slides_data, start=1):
        # 1) Add slide with title & content
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info.title

        # 2) Populate bullets
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for point in slide_info.content:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

        # 3) Auto-generate image if only a prompt is supplied
        if not slide_info.image_path and slide_info.image_prompt:
            img_name = f"slide_{idx}.png"
            result = generate_image_from_prompt.invoke(
                {"prompt": slide_info.image_prompt, "filename": img_name}
            )
            if result.startswith("✅"):
                # Extract the returned path after the arrow
                slide_info.image_path = result.split("→", 1)[-1].strip()

        # 4) Insert the image if available and adjust layout
        if slide_info.image_path and os.path.exists(slide_info.image_path):
            # Adjust text box to make space for the image
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(5.5)
            height = Inches(5.5)
            txBox = slide.shapes.placeholders[1]
            txBox.left, txBox.top, txBox.width, txBox.height = left, top, width, height

            # Add the image to the right of the text
            left_img = Inches(6.0)
            top_img = Inches(1.5)
            width_img = Inches(3.5)
            pic = slide.shapes.add_picture(
                slide_info.image_path,
                left_img, top_img, width=width_img
            )

    # 5) Save & return confirmation
    out_path = os.path.join(AGENT_DIR, output_filename)
    prs.save(out_path)
    return f"✅ Presentation saved to {out_path}"


@tool
def create_docx_document(title: str, content: List[DocxContent], output_filename: str = "document.docx") -> str:
    """
    Creates a Word document with a title, paragraphs, and optional images.
    """
    document = Document()
    document.add_heading(title, level=1)

    for i, item in enumerate(content):
        document.add_paragraph(item.text)
        
        if item.image_path and os.path.exists(item.image_path):
            document.add_picture(item.image_path, width=Inches(5.0))

    output_path = os.path.join(AGENT_DIR, output_filename)
    document.save(output_path)
    return f"✅ Document saved to {output_path}"


